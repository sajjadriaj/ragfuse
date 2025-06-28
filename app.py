# app.py
from flask import Flask, render_template, request, jsonify, redirect, url_for, send_file
import os
import uuid
import re
import hashlib
import logging
import shutil
from datetime import datetime
from werkzeug.utils import secure_filename
from chonkie import NeuralChunker
from duckduckgo_search import DDGS


# Document processing imports
import PyPDF2
import docx
from pptx import Presentation

# ChromaDB import
import chromadb

# Conversation storage import
from conversation.sqlite import SQLiteConversationStorage
from documents.sqlite import SQLiteDocumentStorage
from settings.sqlite import LLMSettingsStorage
from llms.openai_llm import OpenAILLM
from llms.claude_llm import ClaudeLLM
from llms.gemini_llm import GeminiLLM
from llms.ollama_llm import OllamaLLM

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Initialize Flask app
app = Flask(__name__)
app.config["UPLOAD_FOLDER"] = "uploads"
app.config["MAX_CONTENT_LENGTH"] = 50 * 1024 * 1024  # 50MB

# Initialize ChromaDB
try:
    chroma_client = chromadb.PersistentClient(path="/app/chroma_data")
    collection = chroma_client.get_or_create_collection(
        name="knowledge_base", metadata={"hnsw:space": "cosine"}
    )
    logger.info("ChromaDB initialized successfully")
except Exception as e:
    logger.error(f"ChromaDB initialization failed: {e}")
    collection = None

# Constants
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ALLOWED_EXTENSIONS = {"txt", "pdf", "docx", "pptx", "md", "csv", "json"}

# Initialize storage systems
SQLITE_DB_DIR = "/app/sqlite_dbs"

conversation_db: SQLiteConversationStorage = SQLiteConversationStorage(db_path=os.path.join(SQLITE_DB_DIR, "conversations.db"))
conversation_db.init()

# Use a separate DB for documents
DOCUMENT_DB_PATH = os.path.join(SQLITE_DB_DIR, "documents.db")
document_db = SQLiteDocumentStorage(db_path=DOCUMENT_DB_PATH)
document_db.init()

# Initialize LLM settings storage
LLM_SETTINGS_DB_PATH = os.path.join(SQLITE_DB_DIR, "llm_settings.db")
llm_settings_db = LLMSettingsStorage(db_path=LLM_SETTINGS_DB_PATH)
llm_settings_db.init()


def get_llm_instance(llm_provider: str):
    settings = llm_settings_db.get_all_settings()
    if llm_provider == "openai":
        api_key = settings.get("openai_api_key")
        model = settings.get("openai_model", "gpt-3.5-turbo")
        if not api_key: raise ValueError("OpenAI API Key not configured.")
        return OpenAILLM(api_key=api_key, model=model)
    elif llm_provider == "claude":
        api_key = settings.get("claude_api_key")
        model = settings.get("claude_model", "claude-3-sonnet-20240229")
        if not api_key: raise ValueError("Claude API Key not configured.")
        return ClaudeLLM(api_key=api_key, model=model)
    elif llm_provider == "gemini":
        api_key = settings.get("gemini_api_key")
        model = settings.get("gemini_model", "gemini-pro")
        if not api_key: raise ValueError("Gemini API Key not configured.")
        return GeminiLLM(api_key=api_key, model=model)
    elif llm_provider == "ollama":
        endpoint = settings.get("ollama_endpoint", "http://localhost:11434")
        model = settings.get("ollama_model", "llama2")
        return OllamaLLM(endpoint=endpoint, model=model)
    else:
        raise ValueError(f"Unknown LLM provider: {llm_provider}")

def duckduckgo_web_search(query: str) -> dict:
    """Performs a web search using DuckDuckGo Search.
    """
    try:
        # Perform DuckDuckGo search
        # You can adjust the number of results (max_results) as needed
        search_results = DDGS().text(keywords=query, max_results=5)

        formatted_results = []
        for result in search_results:
            formatted_results.append({
                "title": result.get("title"),
                "link": result.get("href"),
                "snippet": result.get("body")
            })
        return {"search_results": formatted_results}

    except Exception as e:
        logger.error(f"DuckDuckGo web search failed: {e}")
        raise RuntimeError(f"Error performing DuckDuckGo web search: {e}")


def allowed_file(filename):
    """Check if file extension is allowed"""
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS


def get_file_hash(file_path):
    """Generate MD5 hash for file"""
    try:
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    except Exception as e:
        logger.error(f"Error generating hash: {e}")
        return None


# Document Processing Functions
def extract_text_from_pdf(file_path):
    """Extract text from PDF file"""
    text = ""
    try:
        with open(file_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                except Exception as e:
                    logger.warning(f"Error extracting text from PDF page: {e}")
                    continue
        logger.info(f"Extracted {len(text)} characters from PDF: {file_path}")
        return text.strip()
    except Exception as e:
        logger.error(f"PDF extraction error for {file_path}: {e}")
        return ""


def extract_text_from_docx(file_path):
    """Extract text from DOCX file"""
    text = ""
    try:
        doc = docx.Document(file_path)
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text += paragraph.text + "\n"

        # Extract from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        text += cell.text + " "
                text += "\n"

        return text.strip()
    except Exception as e:
        logger.error(f"DOCX extraction error: {e}")
        return ""


def extract_text_from_pptx(file_path):
    """Extract text from PPTX file"""
    text = ""
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        logger.error(f"PPTX extraction error: {e}")
        return ""


def extract_text_from_txt(file_path):
    """Extract text from text files"""
    encodings = ["utf-8", "latin-1", "cp1252", "iso-8859-1"]
    for encoding in encodings:
        try:
            with open(file_path, "r", encoding=encoding) as file:
                return file.read().strip()
        except UnicodeDecodeError:
            continue
        except Exception as e:
            logger.error(f"Text extraction error: {e}")
            continue
    return ""


def process_document(file_path, filename):
    """Process document and extract text"""
    try:
        file_extension = filename.rsplit(".", 1)[1].lower()

        # Extract text based on file type
        if file_extension == "pdf":
            text = extract_text_from_pdf(file_path)
        elif file_extension == "docx":
            text = extract_text_from_docx(file_path)
        elif file_extension == "pptx":
            text = extract_text_from_pptx(file_path)
        elif file_extension in ["txt", "md", "csv", "json"]:
            text = extract_text_from_txt(file_path)
        else:
            return ""

        # Clean text
        if text:
            text = re.sub(r"\s+", " ", text)
            text = text.strip()

        return text
    except Exception as e:
        logger.error(f"Document processing error: {e}")
        return ""


# def chunk_text(text, chunk_size=500, overlap_ratio=0.1):
#     """
#     Split text into chunks using sentence boundaries with 10% overlap.
#     """
#     if not text or len(text) <= chunk_size:
#         return [text] if text else []

#     # Split text into sentences (handles ., !, ? and newlines)
#     sentences = re.split(r'(?<=[.!?])\s+|\n+', text)
#     sentences = [s.strip() for s in sentences if s.strip()]
#     chunks = []
#     current_chunk = ""
#     i = 0

#     while i < len(sentences):
#         # Add sentences until chunk_size is reached
#         while i < len(sentences) and len(current_chunk) + len(sentences[i]) + 1 <= chunk_size:
#             current_chunk += (sentences[i] + " ")
#             i += 1

#         current_chunk = current_chunk.strip()
#         if current_chunk:
#             chunks.append(current_chunk)

#         # Overlap: go back by 10% of chunk_size in characters
#         if i < len(sentences):
#             overlap_chars = int(chunk_size * overlap_ratio)
#             # Find where to start the next chunk for overlap
#             overlap_start = max(0, len(current_chunk) - overlap_chars)
#             overlap_text = current_chunk[overlap_start:]
#             # Find the sentence index to resume from
#             j = i
#             overlap_found = False
#             for k in range(i-1, -1, -1):
#                 if overlap_text.startswith(sentences[k]):
#                     j = k
#                     overlap_found = True
#                     break
#             if overlap_found:
#                 i = j
#             else:
#                 # fallback: just go back one sentence
#                 i = max(0, i-1)
#         current_chunk = ""

#     return chunks


# Initialize NeuralChunker
chunker = NeuralChunker(
    model="mirth/chonky_modernbert_base_1",  # Default model
    device_map="cpu",  # Device to run the model on ('cpu', 'cuda', etc.)
    min_characters_per_chunk=10,  # Minimum characters for a chunk
    return_type="chunks",  # Output type
)

def chunk_text(text, chunk_size=10):
    """
    Split text into chunks using NeuralChunker with semantic understanding.
    """
    if not text:
        return []

    if len(text) <= chunk_size:
        return [text]

    # Chunk the text
    chunks = chunker(text)
    return [chunk.text for chunk in chunks]


# Routes
@app.route("/")
def index():
    return redirect(url_for("documents"))


@app.route("/documents")
def documents():
    return render_template("documents.html")


@app.route("/chat")
def chat_page():
    return render_template("chat.html")


@app.route("/settings")
def settings_page():
    return render_template("settings.html")


# API Routes
@app.route("/api/all-documents-and-folders")
def get_all_documents_and_folders():
    """Get all documents and folders for context selection"""
    try:
        data = document_db.get_all_folders_and_files()
        return jsonify(data)
    except Exception as e:
        logger.error(f"Error getting all documents and folders: {e}")
        return jsonify({"error": "Failed to load documents and folders"}), 500


@app.route("/api/folder/<folder_id>")
def get_folder_contents(folder_id):
    """Get folder contents"""
    try:
        folder = document_db.get_folder(folder_id)
        if not folder:
            return jsonify({"error": "Folder not found"}), 404

        subfolders, files = document_db.get_folder_children(folder_id)
        breadcrumb = document_db.build_breadcrumb(folder_id)

        return jsonify(
            {"folder": folder, "breadcrumb": breadcrumb, "contents": subfolders + files}
        )
    except Exception as e:
        logger.error(f"Error getting folder contents: {e}")
        return jsonify({"error": "Failed to load folder"}), 500


@app.route("/api/folder", methods=["POST"])
def create_folder():
    """Create new folder"""
    try:
        data = request.get_json()
        folder_name = data.get("name", "").strip()
        parent_id = data.get("parent_id", "root")

        if not folder_name:
            return jsonify({"error": "Folder name required"}), 400

        if len(folder_name) > 100:
            return jsonify({"error": "Folder name too long"}), 400

        # Check for invalid characters
        invalid_chars = ["/", "\\", ":", "*", "?", '"', "<", ">", "|"]
        if any(char in folder_name for char in invalid_chars):
            return jsonify({"error": "Invalid characters in folder name"}), 400

        # Check for duplicates
        subfolders, _ = document_db.get_folder_children(parent_id)
        if any(f["name"] == folder_name for f in subfolders):
            return jsonify({"error": "Folder already exists"}), 400

        # Create folder
        folder_id = document_db.add_folder(folder_name, parent_id)

        return jsonify(
            {
                "message": "Folder created successfully",
                "folder": {"id": folder_id, "name": folder_name},
            }
        )
    except Exception as e:
        logger.error(f"Error creating folder: {e}")
        return jsonify({"error": "Failed to create folder"}), 500


@app.route("/api/folder/<folder_id>", methods=["DELETE"])
def delete_folder(folder_id):
    """Delete folder"""
    try:
        if folder_id == "root":
            return jsonify({"error": "Cannot delete root folder"}), 400

        if not document_db.get_folder(folder_id):
            return jsonify({"error": "Folder not found"}), 404

        # Also delete from ChromaDB
        if collection:
            try:
                collection.delete(where={"folder_id": folder_id})
            except:
                pass

        document_db.delete_folder(folder_id)

        return jsonify({"message": f"Folder deleted"})
    except Exception as e:
        logger.error(f"Error deleting folder: {e}")
        return jsonify({"error": "Failed to delete folder"}), 500


@app.route("/api/upload", methods=["POST"])
def upload_file():
    try:
        print(f"Upload request received")
        print(f"Files in request: {list(request.files.keys())}")
        print(f"Form data: {dict(request.form)}")

        if not collection:
            return jsonify({"error": "ChromaDB not available"}), 500

        if "file" not in request.files:
            print("No 'file' key in request.files")
            return jsonify({"error": "No file selected"}), 400

        files = request.files.getlist("file")
        folder_id = request.form.get("folder_id", "root")

        print(f"Number of files: {len(files)}")
        print(f"Folder ID: {folder_id}")

        if not files or all(f.filename == "" for f in files):
            print("No files or empty filenames")
            return jsonify({"error": "No files selected"}), 400

        # Check if folder exists in SQLite
        if not document_db.get_folder(folder_id):
            return jsonify({"error": "Folder not found"}), 404

        os.makedirs(app.config["UPLOAD_FOLDER"], exist_ok=True)

        uploaded_files = []
        total_chunks = 0

        for file in files:
            if not file.filename or not allowed_file(file.filename):
                continue

            try:
                filename = secure_filename(file.filename)
                file_extension = filename.rsplit(".", 1)[1].lower()

                # Determine the folder path within UPLOAD_FOLDER
                # If folder_id is 'root', save directly in UPLOAD_FOLDER
                # Otherwise, create a subdirectory for the folder_id
                if folder_id == "root":
                    target_folder_path = app.config["UPLOAD_FOLDER"]
                else:
                    target_folder_path = os.path.join(app.config["UPLOAD_FOLDER"], folder_id)
                    os.makedirs(target_folder_path, exist_ok=True) # Ensure folder exists

                # Define permanent file path
                file_save_path = os.path.join(target_folder_path, filename)
                file.save(file_save_path)

                file_size = os.path.getsize(file_save_path)
                file_hash = get_file_hash(file_save_path)

                if not file_hash:
                    os.remove(file_save_path) # Clean up if hash generation fails
                    continue

                # Process document
                text_content = process_document(file_save_path, filename)
                if not text_content:
                    os.remove(file_save_path) # Clean up if text extraction fails
                    continue

                # Create chunks
                chunks = chunk_text(text_content)
                if not chunks:
                    os.remove(file_save_path) # Clean up if chunking fails
                    continue

                # Add to ChromaDB
                file_id = str(uuid.uuid4())
                chunk_ids = []
                chunk_texts = []
                chunk_metadatas = []

                for i, chunk in enumerate(chunks):
                    chunk_id = f"{file_id}_chunk_{i}"
                    chunk_ids.append(chunk_id)
                    chunk_texts.append(chunk)
                    chunk_metadatas.append(
                        {
                            "file_id": file_id,
                            "filename": filename,
                            "chunk_index": i,
                            "folder_id": folder_id,
                            "file_extension": file_extension,
                            "upload_date": datetime.now().isoformat(),
                        }
                    )

                collection.add(
                    ids=chunk_ids, documents=chunk_texts, metadatas=chunk_metadatas
                )

                # Add to SQLite document storage
                file_info = {
                    "id": file_id,
                    "name": filename,
                    "extension": file_extension,
                    "size": file_size,
                    "hash": file_hash,
                    "folder_id": folder_id,
                    "chunk_count": len(chunks),
                    "created_at": datetime.now().isoformat(),
                    "text_length": len(text_content),
                }
                document_db.add_file(file_info)

                uploaded_files.append(filename)
                total_chunks += len(chunks)

                # No cleanup needed, file is permanently stored

            except Exception as e:
                logger.error(f"Error processing {file.filename}: {e}")
                # Ensure cleanup of the permanently saved file if an error occurs during processing
                if "file_save_path" in locals() and os.path.exists(file_save_path):
                    os.remove(file_save_path)
                continue

        if not uploaded_files:
            return jsonify({"error": "No files processed successfully"}), 400

        return jsonify(
            {
                "message": f"Successfully uploaded {len(uploaded_files)} files",
                "files": uploaded_files,
                "total_chunks": total_chunks,
            }
        )

    except Exception as e:
        logger.error(f"Upload error: {e}")
        return jsonify({"error": "Upload failed"}), 500


@app.route("/api/file/<file_id>", methods=["DELETE"])
def delete_file(file_id):
    """Delete a file"""
    try:
        if not document_db.get_file(file_id):
            return jsonify({"error": "File not found"}), 404

        # Delete from ChromaDB
        if collection:
            try:
                collection.delete(where={"file_id": file_id})
            except:
                pass

        # Get file info to get the filename
        file_info = document_db.get_file(file_id)
        if not file_info:
            return jsonify({"error": "File not found"}), 404

        # Delete from document storage
        document_db.delete_file(file_id)

        # Delete physical file from uploads folder
        file_path = os.path.join(app.config["UPLOAD_FOLDER"], file_info["name"])
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
                logger.info(f"Deleted physical file: {file_path}")
            except Exception as e:
                logger.error(f"Error deleting physical file {file_path}: {e}")
                # Continue with response even if physical file deletion fails

        return jsonify({"message": f"File deleted"})

    except Exception as e:
        logger.error(f"Error deleting file: {e}")
        return jsonify({"error": "Failed to delete file"}), 500


@app.route("/api/file-content/<file_id>")
def get_file_content(file_id):
    """Get the content of a specific file or serve the file directly if it's a PDF."""
    try:
        file_info = document_db.get_file(file_id)
        if not file_info:
            return jsonify({"error": "File not found"}), 404

        # Determine the file's actual path based on its folder_id
        if file_info["folder_id"] == "root":
            file_path = os.path.join(app.config["UPLOAD_FOLDER"], file_info["name"])
        else:
            file_path = os.path.join(app.config["UPLOAD_FOLDER"], file_info["folder_id"], file_info["name"])
        
        # Ensure the file actually exists on disk
        if not os.path.exists(file_path):
            return jsonify({"error": "File not found on disk"}), 404

        file_extension = file_info["extension"].lower()

        if file_extension == "pdf":
            # Serve PDF files directly
            return send_file(file_path, mimetype='application/pdf')
        elif file_extension in ["txt", "md", "csv", "json", "docx", "pptx"]:
            # For other allowed text-based files, extract and return content
            content = process_document(file_path, file_info["name"])
            if not content:
                return jsonify({"error": "Could not extract content from file"}), 500
            return jsonify({"content": content})
        else:
            return jsonify({"error": "Unsupported file type for direct content display"}), 400

    except Exception as e:
        logger.error(f"Error getting file content for {file_id}: {e}")
        return jsonify({"error": "Failed to retrieve file content"}), 500


@app.route("/api/stats")
def get_stats():
    """Get RAGFuse statistics"""
    try:
        stats = document_db.get_stats()

        # Add ChromaDB chunk count
        try:
            collection_count = collection.count() if collection else 0
        except Exception as e:
            logger.warning(f"Error getting ChromaDB count: {e}")
            collection_count = 0

        stats["total_chunks"] = collection_count
        stats["storage_location"] = "/app/chroma_data"

        return jsonify(stats)

    except Exception as e:
        logger.error(f"Stats error: {e}")
        # Return default stats instead of error
        return jsonify(
            {
                "total_files": 0,
                "total_folders": 0,
                "total_chunks": 0,
                "total_size_bytes": 0,
                "file_types": {},
                "storage_location": "./chroma_db",
            }
        )


@app.route("/api/search", methods=["POST"])
def search_knowledge_base():
    """Search the RAGFuse"""
    try:
        if not collection:
            return jsonify({"error": "ChromaDB not available"}), 500

        data = request.get_json()
        query = data.get("query", "").strip()
        folder_id = data.get("folder_id")
        selected_document_ids = data.get("selected_documents", [])
        n_results = data.get("n_results", 5)

        if not query:
            return jsonify({"error": "Query required"}), 400

        # Build where clause
        where_clause = {}
        if selected_document_ids:
            where_clause["file_id"] = {"$in": selected_document_ids}
        elif folder_id and folder_id != "root":
            where_clause["folder_id"] = folder_id

        # Search
        if where_clause:
            results = collection.query(
                query_texts=[query], n_results=n_results, where=where_clause
            )
        else:
            results = collection.query(query_texts=[query], n_results=n_results)

        # Format results
        search_results = []
        if results["documents"][0]:
            for doc, metadata, distance in zip(
                results["documents"][0],
                results["metadatas"][0],
                results["distances"][0],
            ):
                search_results.append(
                    {
                        "content": doc,
                        "filename": metadata["filename"],
                        "chunk_index": metadata["chunk_index"],
                        "folder_id": metadata["folder_id"],
                        "similarity_score": 1 - distance,
                        "upload_date": metadata["upload_date"],
                    }
                )

        return jsonify(
            {
                "results": search_results,
                "query": query,
                "total_results": len(search_results),
            }
        )

    except Exception as e:
        logger.error(f"Search error: {e}")
        return jsonify({"error": "Search failed"}), 500


@app.route("/api/chat", methods=["POST"])
def chat():
    try:
        if not collection:
            return jsonify({"error": "ChromaDB not available"}), 500

        data = request.get_json()
        message = data.get("message", "").strip()
        folder_id = data.get("folder_id")
        conv_id = data.get("conversation_id")  # Get existing conversation ID
        llm_provider = data.get("llm_provider", "openai") # Get selected LLM provider
        web_search_enabled = data.get("web_search_enabled", False)

        if not message:
            return jsonify({"error": "Message required"}), 400

        # Get current conversation history
        messages = []
        if conv_id:
            messages = conversation_db.get_conversation(conv_id)

        # Add user message
        messages.append(
            {
                "role": "user",
                "content": message,
                "timestamp": datetime.now().isoformat(),
            }
        )

        # Prepare context from ChromaDB
        where_clause = {}
        selected_document_ids = data.get("selected_documents", [])

        if selected_document_ids:
            # If specific documents are selected, filter by them
            where_clause["file_id"] = {"$in": selected_document_ids}
        elif folder_id and folder_id != "root":
            # If a folder is selected but no specific documents, filter by folder
            where_clause["folder_id"] = folder_id

        # Prepare context from ChromaDB
        where_clause = {}
        selected_document_ids = data.get("selected_documents", [])

        if selected_document_ids:
            # If specific documents are selected, filter by them
            where_clause["file_id"] = {"$in": selected_document_ids}
        elif folder_id and folder_id != "root":
            # If a folder is selected but no specific documents, filter by folder
            where_clause["folder_id"] = folder_id

        document_context_parts = []
        document_sources = []

        if collection:
            if where_clause:
                results = collection.query(
                    query_texts=[message], n_results=3, where=where_clause
                )
            else:
                # If no folder or specific documents are selected, search all
                results = collection.query(query_texts=[message], n_results=3)

            logger.info(f"ChromaDB query results: {results}")

            if results["documents"] and results["documents"][0]:
                for doc, metadata, distance in zip(
                    results["documents"][0],
                    results["metadatas"][0],
                    results["distances"][0],
                ):
                    document_context_parts.append(doc)
                    document_sources.append(
                        {
                            "filename": metadata["filename"],
                            "chunk_index": metadata["chunk_index"],
                            "similarity": 1 - distance,
                            "type": "document",
                        }
                    )
        logger.info(f"Document sources prepared: {document_sources}")

        web_search_results = []
        web_search_sources = []
        if web_search_enabled:
            try:
                # Perform web search
                web_results = duckduckgo_web_search(query=message)
                if web_results and web_results.get("search_results"):
                    for i, result in enumerate(web_results["search_results"][:3]): # Limit to top 3 web results
                        web_search_results.append(f"Title: {result.get('title')}\nURL: {result.get('link')}\nSnippet: {result.get('snippet')}")
                        web_search_sources.append({
                            "filename": result.get('title', f"Web Result {i+1}"),
                            "url": result.get('link'),
                            "type": "web",
                            "snippet": result.get('snippet')
                        })
                logger.info(f"Web search results: {web_search_results}")
            except Exception as e:
                logger.error(f"Web search error: {e}")
                web_search_results = [f"Error performing web search: {e}"]

        # Combine contexts and sources
        all_context_parts = document_context_parts + web_search_results
        all_sources = document_sources + web_search_sources

        # Generate LLM response
        try:
            llm = get_llm_instance(llm_provider)
            custom_prompt_template = llm_settings_db.get_setting("custom_llm_prompt")

            if all_context_parts:
                context_str = '\n\n'.join(all_context_parts)
                if custom_prompt_template:
                    rag_prompt = custom_prompt_template.replace("{{context}}", context_str).replace("{{query}}", message)
                else:
                    rag_prompt = f"Given the following context:\n\n{context_str}\n\nAnswer the question: {message}"
                bot_response = llm.generate_response(prompt=rag_prompt, context=all_context_parts)
            else:
                if custom_prompt_template:
                    rag_prompt = custom_prompt_template.replace("{{context}}", "").replace("{{query}}", message)
                else:
                    rag_prompt = message
                bot_response = llm.generate_response(prompt=rag_prompt, context=[])
        except ValueError as ve:
            bot_response = f"LLM Configuration Error: {ve}. Please check your settings."
        except Exception as llm_e:
            bot_response = f"Error generating LLM response: {llm_e}"

        # Add bot response
        messages.append(
            {
                "role": "bot",
                "content": bot_response,
                "timestamp": datetime.now().isoformat(),
                "sources": all_sources, # Include all sources
                "context_parts": all_context_parts, # Include all context parts
            }
        )

        # Save conversation
        if not conv_id:
            # Generate new conversation ID
            conv_id = "conv_" + datetime.now().strftime("%Y%m%d_%H%M%S")

        conversation_db.save_conversation(conv_id, messages, selected_documents=data.get("selected_documents", []))

        return jsonify(
            {
                "response": bot_response,
                "sources": all_sources[:3],
                "context_used": len(all_context_parts) > 0,
                "context_parts": all_context_parts, # Include context parts
                "conversation_id": conv_id,
            }
        )

    except Exception as e:
        logger.error(f"Chat error: {e}")
        return jsonify({"error": "Chat failed"}), 500


@app.route("/api/conversations", methods=["GET"])
def list_conversations():
    try:
        convs = conversation_db.get_all_conversations()
        return jsonify({"conversations": convs})
    except Exception as e:
        logger.error(f"Error listing conversations: {e}")
        return jsonify({"conversations": []})  # Return empty list instead of error


@app.route("/api/conversations/<conv_id>", methods=["GET"])
def get_conversation(conv_id):
    try:
        messages = conversation_db.get_conversation(conv_id)
        return jsonify({"conversation_id": conv_id, "messages": messages})
    except Exception as e:
        logger.error(f"Error fetching conversation {conv_id}: {e}")
        return jsonify({"error": "Failed to load conversation"}), 500


@app.route("/api/conversations/<conv_id>", methods=["DELETE"])
def delete_conversation(conv_id):
    try:
        conversation_db.delete_conversation(conv_id)
        return jsonify({"message": "Conversation deleted successfully"})
    except Exception as e:
        logger.error(f"Error deleting conversation {conv_id}: {e}")
        return jsonify({"error": "Failed to delete conversation"}), 500


@app.route("/api/clear-database", methods=["POST"])
def clear_database():
    """Clear all data from ChromaDB and SQLite document storage"""
    global collection  # Move global declaration to the top
    
    try:
        # Clear ChromaDB by deleting and recreating the collection
        if collection:
            try:
                # Delete the entire collection
                chroma_client.delete_collection("knowledge_base")
                
                # Recreate the collection
                collection = chroma_client.get_or_create_collection(
                    name="knowledge_base", metadata={"hnsw:space": "cosine"}
                )
                logger.info("ChromaDB cleared and recreated successfully")
            except Exception as e:
                logger.error(f"Error clearing ChromaDB: {e}")
                return jsonify({"error": "Failed to clear ChromaDB"}), 500
        
        # Clear SQLite document storage (files and folders except root)
        try:
            conn = document_db.conn
            c = conn.cursor()
            
            # Delete all files
            c.execute("DELETE FROM files")
            
            # Delete all folders except root
            c.execute("DELETE FROM folders WHERE id != 'root'")
            
            conn.commit()
            logger.info("SQLite document storage cleared successfully")

            # Clear physical files from uploads folder
            upload_folder = app.config["UPLOAD_FOLDER"]
            for filename in os.listdir(upload_folder):
                file_path = os.path.join(upload_folder, filename)
                try:
                    if os.path.isfile(file_path) or os.path.islink(file_path):
                        os.unlink(file_path)
                    elif os.path.isdir(file_path):
                        shutil.rmtree(file_path)
                except Exception as e:
                    logger.error(f"Error deleting {file_path}: {e}")
            logger.info("Physical files in uploads folder cleared successfully")
            
        except Exception as e:
            logger.error(f"Error clearing SQLite storage: {e}")
            return jsonify({"error": "Failed to clear document storage"}), 500

        return jsonify({
            "message": "Database cleared successfully",
            "cleared": {
                "chromadb": True,
                "documents": True,
                "folders": True
            }
        })

    except Exception as e:
        logger.error(f"Error clearing database: {e}")
        return jsonify({"error": "Failed to clear database"}), 500


@app.route("/api/settings", methods=["GET"])
def get_llm_settings():
    try:
        settings = llm_settings_db.get_all_settings()
        return jsonify(settings)
    except Exception as e:
        logger.error(f"Error getting LLM settings: {e}")
        return jsonify({"error": "Failed to retrieve settings"}), 500


@app.route("/api/settings", methods=["POST"])
def save_llm_settings():
    try:
        data = request.get_json()
        for key, value in data.items():
            llm_settings_db.save_setting(key, value)
        return jsonify({"message": "Settings saved successfully"}), 200
    except Exception as e:
        logger.error(f"Error saving LLM settings: {e}")
        return jsonify({"error": "Failed to save settings"}), 500


# Error handlers
@app.errorhandler(413)
def too_large(e):
    return jsonify({"error": "File too large (max 50MB)"}), 413


@app.errorhandler(404)
def not_found(e):
    return render_template("documents.html"), 404


@app.errorhandler(500)
def server_error(e):
    logger.error(f"Server error: {e}")
    return jsonify({"error": "Internal server error"}), 500


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
